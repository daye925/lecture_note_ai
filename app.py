import streamlit as st
from google import genai
from pptx import Presentation
import PyPDF2
import base64

client = genai.Client(api_key="AIzaSyD0w5Yu8grH7cCIzaWUbAc4ksH8192hZeU")

st.set_page_config(
    page_title="벼락치기 AI 노트",
    page_icon="⚡",
    layout="wide"
)

def img_to_base64(path):
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()

cat_top = img_to_base64("cat_top.png")
cat_treble = img_to_base64("cat_treble.png")
cat_balloon = img_to_base64("cat_balloon.png")
cat_study = img_to_base64("cat_study.png")
banner = img_to_base64("banner.png")
icon_book = img_to_base64("icon_book.png")
icon_star = img_to_base64("icon_star.png")
icon_pencil = img_to_base64("icon_pencil.png")
icon_note = img_to_base64("icon_note.png")
icon_clock = img_to_base64("icon_clock.png")
icon_question = img_to_base64("icon_question.png")
icon_treble = img_to_base64("icon_treble_small.png")

st.markdown(f"""
<style>
header, footer, #MainMenu {{
    visibility: hidden;
}}

.stApp {{
    background: linear-gradient(180deg, #dff3ff 0%, #f7fcff 100%) !important;
    color: #174f91 !important;
}}

.block-container {{
    padding-top: 1rem;
    max-width: 1050px;
}}

section[data-testid="stSidebar"] {{
    background: linear-gradient(180deg, #d4ecff, #f7fcff) !important;
    border-right: 2px solid #b7dcff;
}}

section[data-testid="stSidebar"] * {{
    color: #174f91 !important;
}}

.sidebar-cat {{
    width: 120px;
    display: block;
    margin: 10px auto;
}}

.sidebar-title {{
    font-size: 28px;
    font-weight: 900;
    text-align: center;
    color: #5b9dff;
    text-shadow: 2px 2px white;
    line-height: 1.25;
}}

.hero {{
    display: flex;
    justify-content: space-between;
    align-items: center;
    margin-bottom: 20px;
}}

.hero-title {{
    font-size: 58px;
    font-weight: 900;
    color: #5b9dff;
    text-shadow: 4px 4px 0 white;
}}

.hero-sub {{
    font-size: 19px;
    line-height: 1.8;
    color: #174f91;
}}

.banner {{
    width: 220px;
}}

.main-cat {{
    width: 250px;
}}

.main-card {{
    background: rgba(255,255,255,0.88);
    border: 2px solid #bfe0ff;
    border-radius: 30px;
    padding: 30px;
    box-shadow: 0 12px 28px rgba(93,150,220,0.18);
    margin-bottom: 25px;
}}

.upload-box {{
    border: 2px dashed #8fc3ff;
    border-radius: 24px;
    padding: 35px;
    text-align: center;
    background: rgba(255,255,255,0.55);
    margin-bottom: 14px;
}}

.upload-box h1 {{
    font-size: 42px;
    color: #5b9dff;
}}

.divider {{
    text-align: center;
    color: #4f8edb;
    font-weight: 900;
    margin: 18px 0;
}}

textarea {{
    background: rgba(255,255,255,0.92) !important;
    color: #174f91 !important;
    border: 2px solid #9dccff !important;
    border-radius: 22px !important;
}}

[data-testid="stFileUploader"] {{
    background: rgba(255,255,255,0.72) !important;
    border: 2px dashed #8fc3ff !important;
    border-radius: 22px !important;
    padding: 15px !important;
}}

[data-testid="stFileUploader"] section,
[data-testid="stFileUploader"] div {{
    background: transparent !important;
    color: #174f91 !important;
}}

[data-testid="stFileUploader"] button {{
    background: white !important;
    color: #3f82df !important;
    border: 2px solid #9dccff !important;
    border-radius: 14px !important;
}}

.stButton > button {{
    display: block;
    margin: 20px auto;
    background: linear-gradient(90deg, #5fa7ff, #8dc8ff) !important;
    color: white !important;
    font-size: 21px;
    font-weight: 900;
    border-radius: 999px;
    padding: 15px 45px;
    border: none;
    box-shadow: 0 10px 22px rgba(78,145,220,0.35);
}}

.preview-title {{
    text-align: center;
    font-size: 24px;
    font-weight: 900;
    color: #2f78d7;
    margin: 25px 0;
}}

.feature-card {{
    background: rgba(255,255,255,0.85);
    border: 2px solid #cfe6ff;
    border-radius: 22px;
    padding: 22px;
    min-height: 210px;
    text-align: center;
    box-shadow: 0 8px 18px rgba(93,150,220,0.12);
    margin-bottom: 18px;
}}

.feature-card img.icon {{
    width: 58px;
}}

.feature-card img.cat {{
    width: 48px;
    margin-top: 8px;
}}

.feature-card h4 {{
    color: #2f78d7;
    font-size: 20px;
    margin: 8px 0;
}}

.feature-card p {{
    color: #536f96;
    font-size: 14px;
    line-height: 1.6;
}}

.result-box {{
    background: rgba(255,255,255,0.92);
    border: 2px solid #bfe0ff;
    border-radius: 28px;
    padding: 28px;
    margin-top: 25px;
}}
</style>
""", unsafe_allow_html=True)

with st.sidebar:
    st.markdown(f"""
    <img class="sidebar-cat" src="data:image/png;base64,{cat_top}">
    <div class="sidebar-title">벼락치기<br>AI 노트 🎵</div>
    <hr>
    <p>🏠 홈</p>
    <p>📘 핵심 개념 추출</p>
    <p>⭐ 중요도 분석</p>
    <p>✏️ 예상 문제 생성</p>
    <p>📝 암기 포인트 생성</p>
    <p>⏰ 벼락치기 노트 생성</p>
    <p>❓ 퀴즈 생성</p>
    <hr>
    <p>Made with Gemini AI 💗</p>
    <img style="width:100px;" src="data:image/png;base64,{icon_treble}">
    """, unsafe_allow_html=True)

st.markdown(f"""
<div class="hero">
    <div>
        <img class="banner" src="data:image/png;base64,{banner}">
        <div class="hero-title">벼락치기 AI 노트</div>
        <div class="hero-sub">
        강의 자료를 업로드하면 시험 대비용 학습 자료를<br>
        자동으로 생성해주는 AI 서비스입니다.
        </div>
    </div>
    <img class="main-cat" src="data:image/png;base64,{cat_treble}">
</div>

<div class="main-card">
    <h3>☁️ 수업 자료 업로드</h3>
    <div class="upload-box">
        <h1>☁️</h1>
        <p>PPT, PDF, TXT 파일을 업로드하세요.</p>
        <p>지원 파일 : .pptx, .pdf, .txt</p>
    </div>
</div>
""", unsafe_allow_html=True)

uploaded_file = st.file_uploader(
    "파일 업로드",
    type=["pptx", "pdf", "txt"],
    label_visibility="collapsed"
)

st.markdown('<div class="divider">또는</div>', unsafe_allow_html=True)

st.markdown(f"""
<div class="main-card">
    <h3>✏️ 수업 내용을 직접 입력하세요</h3>
</div>
""", unsafe_allow_html=True)

col1, col2 = st.columns([3, 1])

with col1:
    direct_text = st.text_area(
        "수업 내용 입력",
        height=160,
        placeholder="수업 내용을 여기에 붙여넣으세요.",
        label_visibility="collapsed"
    )

with col2:
    st.markdown(f"""
    <img style="width:150px; margin-top:20px;" src="data:image/png;base64,{cat_study}">
    """, unsafe_allow_html=True)

button_clicked = st.button("⚡ 벼락치기 노트 생성하기 〉")

st.markdown('<div class="preview-title">✨ 생성될 학습 자료 미리보기 ✨</div>', unsafe_allow_html=True)

features = [
    (icon_book, "핵심 개념 추출", "수업에서 가장 중요한 핵심 개념을 정리해줘요."),
    (icon_star, "중요도 분석", "시험에 나올 가능성을 기준으로 중요도를 분석해줘요."),
    (icon_pencil, "예상 문제 생성", "객관식, 단답형, 서술형 문제와 정답을 생성해줘요."),
    (icon_note, "암기 포인트 생성", "시험 직전에 꼭 외워야 할 포인트를 정리해줘요."),
    (icon_clock, "벼락치기 노트 생성", "5분 안에 볼 수 있는 압축 노트를 만들어줘요."),
    (icon_question, "퀴즈 생성", "OX 퀴즈를 통해 복습하고 이해도를 높여줘요."),
]

for i in range(0, 6, 3):
    cols = st.columns(3)
    for col, item in zip(cols, features[i:i+3]):
        icon, title, desc = item
        with col:
            st.markdown(f"""
            <div class="feature-card">
                <img class="icon" src="data:image/png;base64,{icon}">
                <h4>{title}</h4>
                <p>{desc}</p>
                <img class="cat" src="data:image/png;base64,{cat_balloon}">
            </div>
            """, unsafe_allow_html=True)


def extract_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide_num, slide in enumerate(prs.slides, start=1):
        text += f"\\n[슬라이드 {slide_num}]\\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\\n"
    return text


def extract_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page_num, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text()
        if page_text:
            text += f"\\n[페이지 {page_num}]\\n{page_text}\\n"
    return text


def extract_txt(file):
    return file.read().decode("utf-8")


def make_cram_note(text):
    prompt = f"""
    너는 대학생 시험 대비를 도와주는 AI 학습 코치야.
    아래 수업 자료를 바탕으로 '벼락치기 AI 노트'를 만들어줘.

    # 1. 핵심 개념 추출
    # 2. 중요도 분석
    # 3. 예상 문제 생성
    # 4. 암기 포인트 생성
    # 5. 벼락치기 노트 생성
    # 6. 퀴즈 생성

    수업 자료:
    {text}
    """

    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=prompt
    )
    return response.text


lecture_text = ""

if uploaded_file is not None:
    file_name = uploaded_file.name.lower()
    if file_name.endswith(".pptx"):
        lecture_text = extract_pptx(uploaded_file)
    elif file_name.endswith(".pdf"):
        lecture_text = extract_pdf(uploaded_file)
    elif file_name.endswith(".txt"):
        lecture_text = extract_txt(uploaded_file)

if direct_text.strip():
    lecture_text += "\\n" + direct_text

if lecture_text:
    with st.expander("📄 추출된 원문 보기"):
        st.write(lecture_text)

if button_clicked:
    if not lecture_text.strip():
        st.warning("파일을 업로드하거나 텍스트를 입력해주세요.")
    else:
        with st.spinner("AI가 벼락치기 노트를 생성하고 있습니다..."):
            result = make_cram_note(lecture_text)

        st.markdown('<div class="result-box">', unsafe_allow_html=True)
        st.subheader("⚡ 벼락치기 AI 노트 결과")
        st.write(result)

        st.download_button(
            label="📥 결과 다운로드",
            data=result,
            file_name="cram_ai_note.txt",
            mime="text/plain"
        )

        st.markdown("</div>", unsafe_allow_html=True)